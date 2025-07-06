import os
import time
import re
import io
import zipfile
import tempfile
import concurrent.futures # Added for concurrency
import textwrap
import shutil
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
from pptx import Presentation
from docx import Document
import pandas as pd
from difflib import SequenceMatcher
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import Chroma 
from langchain_text_splitters import RecursiveCharacterTextSplitter
import openai 
from typing import List, Tuple, Dict, Any, Generator, Optional
import gc
from tqdm import tqdm
import hashlib 
from collections import defaultdict
import traceback

if not shutil.which("tesseract"):
    raise EnvironmentError(
        "Tesseract OCR is not installed or not in your PATH. "
        "Please install it from https://github.com/tesseract-ocr/tesseract "
        "and ensure the installation directory is in your system PATH."
    )

class FileProcessor:
    """
    Handles the processing of various file types, extracting text content
    and splitting it into manageable chunks for vectorization.
    """
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=50,
            length_function=len,
            is_separator_regex=False,
        )

    def process_file(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process any file type and yield document chunks with metadata."""
        filename = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pdf":
                yield from self._process_pdf(file_path)
            elif ext == ".pptx":
                yield from self._process_pptx(file_path)
            elif ext == ".docx":
                yield from self._process_docx(file_path)
            elif ext in [".xlsx", ".csv"]:
                yield from self._process_excel_csv(file_path)
            elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"]:
                yield from self._process_image(file_path)
            elif ext == ".txt":
                yield from self._process_txt(file_path)
            elif ext == ".zip":
                yield from self._process_zip(file_path)
            else:
                print(f"Skipping unsupported file type: {filename}")
        except Exception as e:
            print(f"Error processing {filename}: {e}")

    def _process_pdf(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        images_dir = os.path.join('static', 'extracted_images')
        os.makedirs(images_dir, exist_ok=True)
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                # Extract images and their positions/captions
                image_info_list = []
                blocks = page.get_text("blocks")
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image['image']
                    ext = base_image['ext']
                    img_hash = hashlib.md5(image_bytes).hexdigest()
                    img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_{doc_hash}_p{page_num+1}_img{img_hash}.{ext}"
                    img_path = os.path.join(images_dir, img_filename)
                    if not os.path.exists(img_path):
                        with open(img_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                    img_url = f"/static/extracted_images/{img_filename}"
                    img_rect = fitz.Rect(img[1], img[2], img[3], img[4]) if len(img) >= 5 else None
                    # Find the closest text block below the image (likely the caption)
                    caption = None
                    if img_rect:
                        min_dist = float('inf')
                        for b in blocks:
                            b_rect = fitz.Rect(b[:4])
                            if b_rect.y0 >= img_rect.y1 and b[4].strip():
                                dist = b_rect.y0 - img_rect.y1
                                if dist < min_dist:
                                    min_dist = dist
                                    caption = b[4].strip()
                    image_info_list.append({
                        'url': img_url,
                        'rect': img_rect,
                        'caption': caption
                    })
                if not text.strip():  # Try OCR if no text is extracted
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = pytesseract.image_to_string(img)
                if text.strip() or image_info_list:
                    chunks = self.text_splitter.split_text(text)
                    chunk_blocks = page.get_text("blocks")
                    chunk_positions = [b for b in chunk_blocks if b[4].strip()]
                    for i, chunk in enumerate(chunks):
                        # Find the closest image(s) to this chunk by y-position
                        chunk_y = None
                        for b in chunk_positions:
                            if chunk[:20] in b[4]:
                                chunk_y = (b[1] + b[3]) / 2
                                break
                        if chunk_y is None:
                            chunk_y = 0
                        # Find the image(s) with the closest y-position
                        min_dist = float('inf')
                        closest_img = None
                        for info in image_info_list:
                            if info['rect']:
                                img_y = (info['rect'].y0 + info['rect'].y1) / 2
                                dist = abs(chunk_y - img_y)
                                if dist < min_dist:
                                    min_dist = dist
                                    closest_img = info
                        image_urls = [closest_img['url']] if closest_img else []
                        image_captions = [closest_img['caption']] if closest_img and closest_img['caption'] else []
                        yield {
                            "text": chunk,
                            "source": os.path.basename(file_path),
                            "page": page_num + 1,
                            "chunk_number": i + 1,
                            "file_type": "pdf",
                            "content_hash": doc_hash,
                            "image_urls": image_urls,
                            "image_captions": image_captions
                        }

    def _process_pptx(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        images_dir = os.path.join('static', 'extracted_images')
        os.makedirs(images_dir, exist_ok=True)
        prs = Presentation(file_path)
        slide_texts = []
        slide_images = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            image_urls = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    slide_text.append(shape.text_frame.text)
                if hasattr(shape, "image") and shape.image:
                    image = shape.image
                    ext = image.ext
                    img_bytes = image.blob
                    img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_{doc_hash}_slide{slide_num+1}_img.{ext}"
                    img_path = os.path.join(images_dir, img_filename)
                    with open(img_path, 'wb') as img_file:
                        img_file.write(img_bytes)
                    img_url = f"/static/extracted_images/{img_filename}"
                    image_urls.append(img_url)
            slide_images.append(image_urls)
            slide_texts.append("\n".join(slide_text))
        if slide_texts:
            combined_text = "\n\n".join(slide_texts)
            chunks = self.text_splitter.split_text(combined_text)
            for i, chunk in enumerate(chunks):
                slide_num = min(i, len(slide_images)-1)
                relevant_image_urls = slide_images[slide_num]
                yield {
                    "text": chunk,
                    "source": os.path.basename(file_path),
                    "chunk_number": i + 1,
                    "file_type": "pptx",
                    "content_hash": doc_hash,
                    "image_urls": relevant_image_urls
                }

    def _process_docx(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        images_dir = os.path.join('static', 'extracted_images')
        os.makedirs(images_dir, exist_ok=True)
        doc = Document(file_path)
        full_text = []
        image_urls = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_bytes = rel.target_part.blob
                ext = rel.target_ref.split('.')[-1]
                img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_{doc_hash}_img.{ext}"
                img_path = os.path.join(images_dir, img_filename)
                with open(img_path, 'wb') as img_file:
                    img_file.write(img_bytes)
                img_url = f"/static/extracted_images/{img_filename}"
                image_urls.append(img_url)
        for p in doc.paragraphs:
            if p.text:
                full_text.append(p.text)
        if full_text:
            combined_text = "\n\n".join(full_text)
            chunks = self.text_splitter.split_text(combined_text)
            for i, chunk in enumerate(chunks):
                yield {
                    "text": chunk,
                    "source": os.path.basename(file_path),
                    "chunk_number": i + 1,
                    "file_type": "docx",
                    "content_hash": doc_hash,
                    "image_urls": image_urls
                }

    def _process_excel_csv(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        images_dir = os.path.join('static', 'extracted_images')
        os.makedirs(images_dir, exist_ok=True)
        # For Excel, images are rare, but let's check for them in xlsx
        image_urls = []
        if ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path)
                for sheet in wb.worksheets:
                    for image in getattr(sheet, '_images', []):
                        img_bytes = image.ref if hasattr(image, 'ref') else image._data()
                        img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_{doc_hash}_{sheet.title}_img.png"
                        img_path = os.path.join(images_dir, img_filename)
                        with open(img_path, 'wb') as img_file:
                            img_file.write(img_bytes)
                        image_urls.append(f"/static/extracted_images/{img_filename}")
            except Exception:
                pass
        if ext == ".xlsx":
            xls = pd.ExcelFile(file_path)
            sheet_names = xls.sheet_names
        else:  # .csv
            sheet_names = [os.path.basename(file_path)]
        for sheet_name in sheet_names:
            if ext == ".xlsx":
                df = pd.read_excel(file_path, sheet_name=sheet_name)
            else:
                df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            if text.strip():
                chunks = self.text_splitter.split_text(text)
                reference_keywords = ["figure", "chart", "graph", "see below", "see image", "as shown"]
                for i, chunk in enumerate(chunks):
                    chunk_lower = chunk.lower()
                    if any(kw in chunk_lower for kw in reference_keywords):
                        relevant_image_urls = image_urls
                    else:
                        relevant_image_urls = []
                    yield {
                        "text": chunk,
                        "source": os.path.basename(file_path),
                        "sheet": sheet_name,
                        "chunk_number": i + 1,
                        "file_type": ext[1:],
                        "content_hash": doc_hash,
                        "image_urls": relevant_image_urls
                    }

    def _process_image(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        if text.strip():
            chunks = self.text_splitter.split_text(text)
            for i, chunk in enumerate(chunks):
                yield {
                    "text": chunk,
                    "source": os.path.basename(file_path),
                    "chunk_number": i + 1,
                    "file_type": os.path.splitext(file_path)[1][1:],
                    "content_hash": doc_hash
                }

    def _process_txt(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        doc_hash = self._get_file_hash(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        if text.strip():
            chunks = self.text_splitter.split_text(text)
            for i, chunk in enumerate(chunks):
                yield {
                    "text": chunk,
                    "source": os.path.basename(file_path),
                    "chunk_number": i + 1,
                    "file_type": "txt",
                    "content_hash": doc_hash
                }

    def _process_zip(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Processes files within a zip archive."""
        doc_hash = self._get_file_hash(file_path)
        with zipfile.ZipFile(file_path, 'r') as zf:
            with tempfile.TemporaryDirectory() as temp_dir:
                zf.extractall(temp_dir)
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        full_path = os.path.join(root, file)
                        for chunk_data in self.process_file(full_path):
                            chunk_data["source"] = os.path.basename(file_path)
                            chunk_data["original_file_in_zip"] = os.path.basename(full_path)
                            chunk_data["content_hash_zip"] = doc_hash
                            yield chunk_data

    def _get_file_hash(self, file_path: str) -> str:
        """Generates an MD5 hash of the file content."""
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            while chunk := f.read(8192):
                hasher.update(chunk)
        return hasher.hexdigest()


class Chatbot:
    """
    A document chat assistant that processes various document types, stores their
    content in a vector database, and answers user queries using an LLM.
    """
    def __init__(self, data_directory: str = "data", persist_directory: str = "chroma_db"):
        self.data_directory = data_directory
        self.persist_directory = persist_directory
        self.file_processor = FileProcessor()
        
        # --- OpenAI API Key and Client Setup ---
        self.openai_api_key = os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError("OPENAI_API_KEY environment variable not set. Please set it before running.")
        
        # Consider a lower timeout for faster error detection if API is unresponsive
        self.client = openai.OpenAI(api_key=self.openai_api_key, timeout=60.0) # Reduced timeout for faster feedback
        self.llm_model = "gpt-4o-mini" # A fast and cost-effective model

        self.embedding_function = OpenAIEmbeddings(
            openai_api_key=self.openai_api_key,
            model="text-embedding-3-small", # Efficient embedding model
            chunk_size=500 # Ensure chunk_size matches or is a multiple of text splitter's for efficient batching
        )
        # --- End Setup ---

        self.response_cache = {}
        self.file_content_hashes = {}

        os.makedirs(self.data_directory, exist_ok=True)
        os.makedirs(self.persist_directory, exist_ok=True)

        self.vector_store = self.get_vector_store()
        self.populate_database()

    def get_vector_store(self) -> Chroma:
        """Initializes or loads the Chroma vector store."""
        if os.path.exists(self.persist_directory) and os.listdir(self.persist_directory):
            print(f"Loading vector store from {self.persist_directory}")
            return Chroma(
                persist_directory=self.persist_directory,
                embedding_function=self.embedding_function
            )
        else:
            print(f"Creating new vector store in {self.persist_directory}")
            return Chroma(
                embedding_function=self.embedding_function,
                persist_directory=self.persist_directory
            )

    def populate_database(self):
        """Populates the vector database with documents from the data directory."""
        print("Populating database with documents from the 'data' directory...")
        existing_hashes = self._get_existing_document_hashes()
        
        files_to_process = [f for f in os.listdir(self.data_directory) if os.path.isfile(os.path.join(self.data_directory, f))]
        
        # Use ThreadPoolExecutor for concurrent file processing (CPU-bound text extraction)
        documents_to_add_flat = [] # Flattened list of all chunks from new files
        processed_file_hashes = set()

        with concurrent.futures.ThreadPoolExecutor() as executor:
            future_to_file = {
                executor.submit(self._process_single_file_for_populate, file_path, existing_hashes): file_path
                for file_path in [os.path.join(self.data_directory, f) for f in files_to_process]
            }
            
            for future in tqdm(concurrent.futures.as_completed(future_to_file), total=len(files_to_process), desc="Processing files"):
                file_path = future_to_file[future]
                try:
                    file_hash, chunks_data = future.result()
                    if file_hash and chunks_data:
                        documents_to_add_flat.extend(chunks_data)
                        self.file_content_hashes[file_hash] = os.path.basename(file_path)
                        processed_file_hashes.add(file_hash)
                    elif file_hash: # Means file was skipped due to being identical
                         self.file_content_hashes[file_hash] = os.path.basename(file_path) # Ensure it's in the map
                except Exception as e:
                    print(f"Error processing {os.path.basename(file_path)} during population: {e}")

        if documents_to_add_flat:
            print(f"\nAdding {len(documents_to_add_flat)} new text chunks to the vector store...")
            # Langchain's add_texts method already handles batching embeddings efficiently
            self.vector_store.add_texts(
                texts=[doc["text"] for doc in documents_to_add_flat],
                metadatas=[{k: v for k, v in doc.items() if k != "text"} for doc in documents_to_add_flat]
            )
            self.vector_store.persist()
            print(f"Database population complete. Added {len(processed_file_hashes)} new files.")
        else:
            print("No new documents to add or all files already processed.")
        gc.collect()

    def _process_single_file_for_populate(self, file_path: str, existing_hashes: set) -> Tuple[Optional[str], Optional[List[Dict[str, Any]]]]:
        """Helper function to process a single file for populate_database."""
        if not os.path.exists(file_path):
            return None, None
        
        file_hash = self.file_processor._get_file_hash(file_path)
        if file_hash in existing_hashes:
            return file_hash, None # Indicate file was skipped
        
        print(f"    Adding {os.path.basename(file_path)} to database...")
        chunks_data = list(self.file_processor.process_file(file_path))
        
        if not chunks_data:
            print(f"    No extractable text found in {os.path.basename(file_path)}.")
            return file_hash, None # Indicate no text was extracted
        
        return file_hash, chunks_data


    def _get_existing_document_hashes(self) -> set:
        """Retrieves content hashes of documents already in the Chroma database."""
        try:
            # Fetch existing metadata from the vector store
            all_data = self.vector_store.get(include=["metadatas"])
            all_metadatas = all_data.get("metadatas", [])
            
            # Populate the file_content_hashes map for efficient lookups
            for metadata in all_metadatas:
                if 'content_hash' in metadata and 'source' in metadata:
                    self.file_content_hashes[metadata['content_hash']] = metadata['source']
                # Handle old zip content hash if present
                if 'content_hash_zip' in metadata and 'source' in metadata:
                    self.file_content_hashes[metadata['content_hash_zip']] = metadata['source']

            # Return a set of all unique content hashes (both file and zip hashes)
            return {m.get('content_hash') for m in all_metadatas if 'content_hash' in m} | \
                   {m.get('content_hash_zip') for m in all_metadatas if 'content_hash_zip' in m}
        except Exception as e:
            if "does not exist" in str(e) or "No existing DB" in str(e):
                return set()
            print(f"Warning: Could not retrieve existing document hashes: {e}")
            return set()

    def add_document(self, file_path: str):
        """Adds a single document to the database."""
        if not os.path.exists(file_path):
            print(f"Error: File not found at {file_path}")
            return

        filename = os.path.basename(file_path)
        destination_path = os.path.join(self.data_directory, filename)
        new_file_hash = self.file_processor._get_file_hash(file_path)

        if new_file_hash in self.file_content_hashes:
            print(f"File '{filename}' is identical to an existing document ('{self.file_content_hashes[new_file_hash]}'). Skipping.")
            return

        print(f"Copying '{filename}' to '{self.data_directory}'...")
        try:
            shutil.copy(file_path, destination_path)
            print(f"Processing '{filename}'...")
            documents_to_add = list(self.file_processor.process_file(destination_path))

            if documents_to_add:
                # Add chunks in batches for efficiency
                self.vector_store.add_texts(
                    texts=[doc["text"] for doc in documents_to_add],
                    metadatas=[{k: v for k, v in doc.items() if k != "text"} for doc in documents_to_add]
                )
                self.vector_store.persist()
                self.file_content_hashes[new_file_hash] = filename
                print(f"Successfully added '{filename}' to the database.")
            else:
                print(f"No extractable text found in '{filename}'.")
        except Exception as e:
            print(f"Failed to add document '{filename}': {e}")
        gc.collect()

    def clear_database(self):
        """Clears the Chroma database and deletes all documents from the data directory."""
        confirmation = input("Are you sure you want to clear the entire database and delete all files? (yes/no): ").lower()
        if confirmation == 'yes':
            if os.path.exists(self.persist_directory):
                shutil.rmtree(self.persist_directory)
            
            if os.path.exists(self.data_directory):
                for filename in os.listdir(self.data_directory):
                    file_path = os.path.join(self.data_directory, filename)
                    if os.path.isfile(file_path):
                        os.remove(file_path)

            self.vector_store = self.get_vector_store()
            self.file_content_hashes = {}
            print("Database and data directory have been cleared.")
            gc.collect()
        else:
            print("Operation cancelled.")

    def reload_database(self):
        """Clears the database and repopulates it with current files in the data directory."""
        confirmation = input("Are you sure you want to reload the database? (yes/no): ").lower()
        if confirmation == 'yes':
            print("Reloading database...")
            if os.path.exists(self.persist_directory):
                shutil.rmtree(self.persist_directory)
            self.vector_store = self.get_vector_store()
            self.file_content_hashes = {}
            self.populate_database()
            print("Database reloaded successfully.")
            gc.collect()
        else:
            print("Operation cancelled.")

    def query_vector_store(self, query_text: str, k: int = 5) -> List[Dict[str, Any]]:
        """Queries the vector store for relevant document chunks."""
        # This method is already using similarity_search_with_score which is efficient for retrieval
        results = self.vector_store.similarity_search_with_score(query_text, k=k)
        return [{"text": doc.page_content, "metadata": doc.metadata, "score": score} for doc, score in results]

    def _get_relevant_sources(self, response_text: str, context_chunks: List[Dict[str, Any]]) -> List[str]:
        """Identifies unique source filenames that contributed to the response."""
        relevant_sources = set()
        
        # Use concurrent.futures for faster similarity checks if context_chunks is large
        with concurrent.futures.ThreadPoolExecutor() as executor:
            # Pass the original response_text to the helper
            futures = {executor.submit(self._check_chunk_relevance, response_text, chunk): chunk for chunk in context_chunks}
            for future in concurrent.futures.as_completed(futures):
                source = future.result()
                if source:
                    relevant_sources.add(source)
                
        return sorted(list(relevant_sources))

    def _check_chunk_relevance(self, response_text: str, chunk: Dict[str, Any]) -> Optional[str]:
        """
        Helper to check if a single chunk is relevant to the response. It uses a multi-pronged
        approach: direct sentence quotation, high keyword overlap per sentence, and a fallback
        for highly relevant chunks.
        """
        source_filename = chunk["metadata"].get("source")
        if not source_filename:
            return None

        # --- Helper function for cleaning text ---
        def _normalize_text(text: str) -> str:
            # Keep alphanumeric characters and periods (for numbers like 99.5), then convert to lowercase.
            processed_text = re.sub(r'[^\w\s\.]', '', text).lower()
            # Collapse multiple whitespace characters into a single space
            return re.sub(r'\s+', ' ', processed_text).strip()

        # --- Helper function for extracting keywords ---
        def _get_keywords(text: str) -> set:
            STOP_WORDS = set([
                "a", "about", "above", "after", "again", "against", "all", "am", "an", "and", "any", "are", "as", "at",
                "be", "because", "been", "before", "being", "below", "between", "both", "but", "by", "can", "did", "do",
                "does", "doing", "down", "during", "each", "few", "for", "from", "further", "had", "has", "have", "having",
                "he", "her", "here", "hers", "herself", "him", "himself", "his", "how", "i", "if", "in", "into", "is", "it",
                "its", "itself", "just", "me", "more", "most", "my", "myself", "no", "nor", "not", "now", "of", "off", "on",
                "once", "only", "or", "other", "our", "ours", "ourselves", "out", "over", "own", "s", "same", "she", "should",
                "so", "some", "such", "t", "than", "that", "the", "their", "theirs", "them", "themselves", "then", "there",
                "these", "they", "this", "those", "through", "to", "too", "under", "until", "up", "very", "was", "we", "were",
                "what", "when", "where", "which", "while", "who", "whom", "why", "will", "with", "you", "your", "yours",
                "yourself", "yourselves"
            ])
            normalized_text = _normalize_text(text)
            words = normalized_text.split()
            return {word for word in words if word not in STOP_WORDS}

        # --- Pre-computation for the response ---
        normalized_response = _normalize_text(response_text)
        response_keywords = _get_keywords(response_text)
        
        # Split chunk into sentences. This regex is basic but effective for many cases.
        chunk_sentences = re.split(r'(?<=[.!?])\s+', chunk["text"])

        # --- Iterate through each sentence of the chunk ---
        for sentence in chunk_sentences:
            if len(sentence.split()) < 3: # Ignore very short, likely meaningless sentences
                continue

            # --- Strategy 1: Direct Substring Match (for quotations) ---
            normalized_sentence = _normalize_text(sentence)
            if normalized_sentence and normalized_sentence in normalized_response:
                return source_filename

            # --- Strategy 2: High Keyword Overlap (for rephrasing) ---
            sentence_keywords = _get_keywords(sentence)
            if not sentence_keywords:
                continue
            
            common_keywords = sentence_keywords.intersection(response_keywords)
            # If a high percentage of a sentence's keywords are in the response, it's a match.
            overlap_percentage = (len(common_keywords) / len(sentence_keywords)) * 100
            
            if overlap_percentage > 70:
                return source_filename

        # --- Strategy 3: Fallback for Highly Relevant Chunks (Vector Similarity) ---
        # This checks the whole chunk against the response, useful if the LLM synthesizes
        # info from the entire chunk rather than a single sentence.
        chunk_score = chunk.get("score")
        if chunk_score is not None and chunk_score < 0.5: # 0.5 is a heuristic for "highly relevant"
            chunk_keywords = _get_keywords(chunk["text"])
            if not chunk_keywords:
                return None
                
            common_keywords = chunk_keywords.intersection(response_keywords)
            # If the chunk is highly relevant and shares a decent number of keywords, cite it.
            # A lower threshold is used here because the chunk might be longer than the response.
            overlap_percentage = (len(common_keywords) / len(chunk_keywords)) * 100
            if overlap_percentage > 30:
                return source_filename
                
        return None

    def list_documents(self) -> str:
        """Lists all the files currently present in the 'data' directory."""
        files = [f for f in os.listdir(self.data_directory) if os.path.isfile(os.path.join(self.data_directory, f))]
        if not files:
            return "There are no documents currently loaded in the database."
        else:
            sorted_files = sorted(files, key=lambda s: s.lower())
            return "Documents currently in the database:\n" + "\n".join([f"- {file}" for file in sorted_files])

    def generate_response(self, user_query: str):
        """
        Generates a response to the user query using the LLM, incorporating
        relevant document context. Returns both the LLM summary and all relevant chunks with their sources.
        """
        if user_query in self.response_cache:
            return self.response_cache[user_query]
        
        # Get more chunks, then filter/merge
        raw_chunks = self.query_vector_store(user_query, k=20)
        # Sort by score (lower is better)
        sorted_chunks = sorted(raw_chunks, key=lambda x: x.get('score', 1e9))
        # Merge consecutive chunks from the same source
        merged_chunks = []
        prev = None
        for chunk in sorted_chunks:
            meta = chunk['metadata']
            if prev and prev['metadata'].get('source') == meta.get('source') and prev['metadata'].get('page') == meta.get('page'):
                prev['text'] += '\n' + chunk['text']
                # Merge images/captions
                prev['metadata']['image_urls'] = list(set(prev['metadata'].get('image_urls', []) + meta.get('image_urls', [])))
                prev['metadata']['image_captions'] = list(set(prev['metadata'].get('image_captions', []) + meta.get('image_captions', [])))
            else:
                merged_chunks.append(chunk)
                prev = chunk
        # Only show top 5
        context_chunks = merged_chunks[:5]
        
        system_message_content = """
        You are a helpful AI assistant. Use the provided document context to answer the user's query accurately.
        If the information is not in the context, state that the documents do not contain the answer.
        Do not make up information. Base your answers strictly on the provided context.
        """
        
        context_text = "\n\n".join([f"Source: {chunk['metadata'].get('source', 'Unknown')}\nContent: {chunk['text']}" for chunk in context_chunks])
        user_message_content = f"User Query: {user_query}\n\nDocument Context:\n{context_text}"
        
        messages = [
            {"role": "system", "content": system_message_content},
            {"role": "user", "content": user_message_content},
        ]
        
        try:
            response_data = self.client.chat.completions.create(
                model=self.llm_model,
                messages=messages
            )
            full_response = response_data.choices[0].message.content

            # Prepare all relevant chunks with their sources
            answer_chunks = []
            for chunk in context_chunks:
                answer_chunks.append({
                    "text": chunk["text"],
                    "source": chunk["metadata"].get("source", "Unknown"),
                    "metadata": chunk["metadata"]
                })
            # Unique sources for backward compatibility
            relevant_sources = list({chunk["source"] for chunk in answer_chunks})

            result = {
                "summary": full_response,
                "chunks": answer_chunks,
                "sources": relevant_sources
            }
            self.response_cache[user_query] = result
            return result
        except Exception as e:
            return {"summary": f"An error occurred while communicating with OpenAI: {e}", "chunks": [], "sources": []}

    def analyze_file_similarity(self, new_file_path: str, similarity_threshold: float = 0.65):
        """
        Analyzes the similarity of a new file against existing documents.
        Returns a string summary of the analysis result.
        """
        if not os.path.exists(new_file_path):
            return f"Error: File not found at {new_file_path}"

        new_file_hash = self.file_processor._get_file_hash(new_file_path)
        if new_file_hash in self.file_content_hashes:
            return f"The file '{os.path.basename(new_file_path)}' is identical to an existing document ('{self.file_content_hashes[new_file_hash]}')."

        new_file_chunks = list(self.file_processor.process_file(new_file_path))
        if not new_file_chunks:
            return f"No extractable text found in '{os.path.basename(new_file_path)}'."

        total_new_chunks = len(new_file_chunks)
        highly_similar_matches = defaultdict(int)

        # Use ThreadPoolExecutor for concurrent similarity checks for chunks
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future_to_chunk = {
                executor.submit(self._find_similar_chunk_in_db, chunk["text"], similarity_threshold): chunk
                for chunk in new_file_chunks
            }
            for future in concurrent.futures.as_completed(future_to_chunk):
                new_chunk = future_to_chunk[future]
                try:
                    similar_doc_source = future.result()
                    if similar_doc_source and similar_doc_source != os.path.basename(new_file_path):
                        highly_similar_matches[similar_doc_source] += 1
                except Exception as e:
                    pass  # Optionally log error

        if highly_similar_matches:
            result_lines = [f"'{os.path.basename(new_file_path)}' shows high similarity to the following documents:"]
            for filename, count in highly_similar_matches.items():
                percentage = (count / total_new_chunks) * 100
                result_lines.append(f"- {filename}: {percentage:.2f}% of chunks are highly similar.")
            return "\n".join(result_lines)
        else:
            return f"'{os.path.basename(new_file_path)}' does not show significant content overlap with existing documents."

    def _find_similar_chunk_in_db(self, text: str, similarity_threshold: float) -> Optional[str]:
        """Helper to find if a single text chunk has a highly similar match in DB."""
        similar_docs = self.vector_store.similarity_search_with_score(text, k=1)
        if similar_docs:
            doc, score = similar_docs[0]
            if score < similarity_threshold: 
                return doc.metadata.get('source')
        return None

    def chat_interface(self):
        """Provides an interactive command-line chat interface."""
        print("\n--- Document Chat Assistant ---")
        print("Type /help for commands or ask naturally.")
        
        command_phrases = {
            "add": [
                "add a new document", "add this file", "include this document",
                "upload a file", "add document", "add this pdf", "add to database",
                "ingest this file", "process this document", "store this file",
                "add file to system", "add new file", "include in database",
                "upload document", "add this doc", "add to collection",
                "add this to knowledge base", "add this to the system",
                "please add this file", "can you add this document",
                "add the file", "add a document", "add file",
                "add this to documents", "add to doc store", "add to vector store",
                "incorporate this file", "include this pdf", "process this file",
                "add this content", "store document", "include this report",
                "add the document", "upload this file", "add this to the database",
                "please include this document", "can you process this file",
                "add this to your knowledge", "store this document",
                "add this material", "include this content", "add this resource",
                "upload the file", "add this to the system", "incorporate document",
                "add this to collection", "process the file", "store the document"
            ],
            "clear": [
                "clear the database", "reset the database", "delete all documents",
                "remove all files", "erase the database", "wipe all data",
                "clean the database", "start fresh", "clear everything",
                "purge database", "delete everything", "remove all content",
                "clear all documents", "reset system", "empty the database",
                "wipe the database", "clear knowledge base", "delete all files",
                "remove all data", "erase everything", "clean slate",
                "purge all documents", "delete database", "clear vector store",
                "reset to empty", "clean out the database", "remove everything",
                "delete all content", "erase all files", "reset document store",
                "flush the database", "clear all data", "remove documents",
                "delete the knowledge base", "erase the vector store",
                "reset the system", "clean the knowledge base", "purge files",
                "delete stored documents", "remove all stored files",
                "erase document collection", "clear all", "reset everything",
                "clean out everything", "delete all", "remove all",
                "purge all", "flush everything", "erase all"
            ],
            "reload": [
                "reload the database", "refresh the documents", "reprocess the files",
                "reload documents", "update database", "refresh knowledge base",
                "reindex documents", "reprocess all files", "reload the system",
                "refresh database", "reinitialize database", "reload files",
                "update the collection", "reprocess documents", "reindex files",
                "reload from data directory", "reprocess all", "refresh documents",
                "reload knowledge base", "update vector store", "reinitialize system",
                "reindex database", "refresh the system", "reload all documents",
                "reprocess the database", "refresh the collection", "reindex all",
                "update the knowledge base", "reload the document store",
                "refresh vector store", "reprocess the collection", "reinitialize",
                "reindex the system", "update stored documents", "refresh all files",
                "reload everything", "reprocess knowledge base", "reindex documents",
                "update all documents", "refresh the database", "reload content",
                "reprocess the vector store", "reindex the knowledge base",
                "update the system", "refresh everything", "reload data",
                "reprocess the content", "reindex everything", "update all"
            ],
            "analyze": [
                "analyze this file", "analyze that file", "analyze the file",
                "please analyze", "pls analyze", "check similarity",
                "compare file", "file similarity", "how similar is",
                "analyze similarity", "analysis for", "run analysis on",
                "process analysis for", "examine similarity", "evaluate similarity",
                "check content match", "find similar documents", "compare to database",
                "document similarity", "how similar to others", "content similarity",
                "analyze for duplicates", "find matching documents", "check for similar content",
                "analyze document similarity", "compare document", "similarity analysis",
                "how similar is this file", "find similar content", "check document match",
                "evaluate file similarity", "examine content match", "compare with database",
                "determine similarity", "check for duplicates", "find related documents",
                "assess similarity", "measure content similarity", "compare to other files",
                "check how similar", "find matching content", "analyze against database",
                "evaluate document match", "examine file similarity", "compare content",
                "check similarity score", "find duplicate content", "assess document similarity",
                "measure file similarity", "determine content match"
            ],
            "source": [
                "list documents", "show files", "what documents are loaded",
                "list sources", "show available documents", "list all files",
                "what files are in the database", "show stored documents",
                "display documents", "list loaded files", "show database contents",
                "what docs are available", "list all documents", "show source files",
                "display file list", "list knowledge base", "show stored files",
                "what documents do you have", "list available files", "show document collection",
                "display document list", "list all sources", "show all documents",
                "what files are stored", "list database files", "show me the documents",
                "display available files", "list stored documents", "show what's loaded",
                "what content do you have", "list files in system", "show document inventory",
                "display stored files", "list available content", "show documents in database",
                "what is in the knowledge base", "list stored content", "show me the files",
                "display knowledge base", "list all stored documents", "show the document list",
                "what do you have stored", "list loaded documents", "show available content",
                "display the document collection", "list all files in system",
                "show me all documents", "display all files", "list everything stored"
            ],
            "exit": [
                "exit the program", "quit", "close the application", "bye",
                "exit application", "end session", "close program", "terminate",
                "goodbye", "shutdown", "leave", "finish", "stop", "end",
                "exit now", "quit now", "close now", "bye bye", "see you later",
                "that's all", "i'm done", "end conversation", "terminate program",
                "close chatbot", "exit system", "quit application", "shut down",
                "log out", "sign off", "terminate session", "close down",
                "end the program", "quit the application", "close the session",
                "exit from here", "goodbye for now", "i'm finished", "end now",
                "terminate now", "close interface", "exit chat", "quit system",
                "close down program", "end chat session", "finish conversation",
                "terminate application", "close assistant", "exit now please",
                "quit the program", "close the chatbot", "end this session"
            ]
        }

        # Initial population of the database (already in __init__, but useful to keep context)
        # This part runs when Chatbot is initialized in main.py, so no need to repeat here.

        while True:
            try:
                user_input = input("\nYou: ").strip()
                if not user_input:
                    continue

                lower_input = user_input.lower()
                
                if lower_input == "/help":
                    print("\nCommands:\n"
                          "  /add <file_path>      - Add a new document\n"
                          "  /clear                - Clear the database\n"
                          "  /reload               - Reload the database\n"
                          "  /analyze <file_path>  - Analyze file similarity\n"
                          "  /source               - List documents\n"
                          "  /exit                 - Exit the chat\n"
                          "\nNatural language examples:\n"
                          "  \"Add this document: /path/to/file.pdf\"\n"
                          "  \"Please analyze this file for similarity\"\n"
                          "  \"What files are in the database?\"\n"
                          "  \"Reset the knowledge base\"")
                    continue
                    
                # Handle explicit slash commands first
                if lower_input.startswith("/add "):
                    file_path = user_input[len("/add "):].strip()
                    self.add_document(file_path)
                    continue
                elif lower_input == "/clear":
                    self.clear_database()
                    continue
                elif lower_input == "/reload":
                    self.reload_database()
                    continue
                elif lower_input.startswith("/analyze "):
                    file_path = user_input[len("/analyze "):].strip()
                    print(self.analyze_file_similarity(file_path))
                    continue
                elif lower_input == "/source":
                    print(self.list_documents())
                    continue
                elif lower_input == "/exit":
                    print("Exiting...")
                    break

                # Handle natural language requests
                command_detected = None
                for command, phrases in command_phrases.items():
                    for phrase in phrases:
                        if phrase in lower_input:
                            command_detected = command
                            break
                    if command_detected:
                        break
                
                file_path = None
                if command_detected in ["add", "analyze"]:
                    quoted_match = re.search(r'[\'"](.+?)[\'"]', user_input)
                    if quoted_match:
                        file_path = quoted_match.group(1)
                    
                    if not file_path:
                        path_match = re.search(r'(\S+\.\w{3,4})\b', user_input) 
                        if path_match:
                            file_path = path_match.group(1)
                
                if command_detected == "add":
                    if file_path:
                        self.add_document(file_path)
                    else:
                        print("Assistant: Please specify a file path to add. Example: 'add document \"C:\\path\\to\\file.pdf\"'")
                elif command_detected == "clear":
                    self.clear_database()
                elif command_detected == "reload":
                    self.reload_database()
                elif command_detected == "analyze":
                    if file_path:
                        print(self.analyze_file_similarity(file_path))
                    else:
                        print("Assistant: Please specify a file path for similarity analysis. Example: 'analyze this file \"C:\\path\\to\\file.docx\"'")
                elif command_detected == "source":
                    print(self.list_documents())
                elif command_detected == "exit":
                    print("Exiting...")
                    break
                else: # Default to RAG if no command is detected
                    start_time = time.perf_counter()
                    response = self.generate_response(user_input)
                    end_time = time.perf_counter()
                    
                    print("\nAssistant:")
                    print(textwrap.fill(response["summary"], width=90, subsequent_indent="  "))
                    
                    source_info = f"Sources: {', '.join(response['sources'])}" if response['sources'] else "No specific sources cited."
                    print(f"\n[{end_time - start_time:.2f}s | {source_info}]")

            except Exception as e:
                print(f"\nAn unexpected error occurred in the chat interface: {e}")

    def delete_document(self, filename: str) -> str:
        """Deletes a specific document and its vector data from the database by filename. Logs locked files if deletion fails."""
        import traceback
        file_path = os.path.join(self.data_directory, filename)
        if not os.path.isfile(file_path):
            return f"File '{filename}' not found in the database."
        try:
            try:
                os.remove(file_path)
            except PermissionError as pe:
                return f"Failed to delete '{filename}': File is in use or locked by another process. Please close any programs using this file and try again."
            except Exception as e:
                tb = traceback.format_exc()
                return f"Failed to delete '{filename}': {e}\nTraceback: {tb}"
            # Release the vector store reference before deleting the directory
            self.vector_store = None
            import gc
            gc.collect()
            if os.path.exists(self.persist_directory):
                try:
                    shutil.rmtree(self.persist_directory)
                except PermissionError:
                    return f"Failed to delete vector store directory '{self.persist_directory}': Directory is in use or locked by another process. Please close any programs using files in this directory and try again."
                except Exception as e:
                    tb = traceback.format_exc()
                    return f"Failed to delete vector store directory '{self.persist_directory}': {e}\nTraceback: {tb}"
            self.vector_store = self.get_vector_store()
            self.file_content_hashes = {}
            self.populate_database()
            return f"File '{filename}' and its data have been deleted."
        except Exception as e:
            tb = traceback.format_exc()
            return f"Failed to delete '{filename}': {e}\nTraceback: {tb}"
